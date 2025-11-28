import math
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# 1. THEME CONFIGURATION (深空玻璃态设计系统)
# ==============================================================================
class Theme:
    # 尺寸 (16:9 Wide)
    WIDTH = Inches(13.333)
    HEIGHT = Inches(7.5)
    
    # 颜色定义 (Dark Mode + Neon Accents)
    COLOR_BG_DEEP = RGBColor(13, 17, 23)       # GitHub Darker Dimmed
    COLOR_BG_CARD = RGBColor(30, 35, 45)       # Card Background
    COLOR_TEXT_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_GREY = RGBColor(160, 170, 190)
    
    # 角色主题色 (Role-Coding)
    COLOR_FARMER = RGBColor(56, 239, 125)      # Neon Green
    COLOR_BUYER = RGBColor(56, 189, 248)       # Sky Blue
    COLOR_BANK = RGBColor(255, 215, 0)         # Gold
    COLOR_EXPERT = RGBColor(192, 132, 252)     # Purple
    COLOR_ADMIN = RGBColor(244, 63, 94)        # Rose Red
    COLOR_TECH = RGBColor(0, 214, 194)         # React Cyan

    # 字体
    FONT_EN = "Segoe UI"    # 英文无衬线
    FONT_CN = "Microsoft YaHei UI" # 中文无衬线

# ==============================================================================
# 2. SLIDE GENERATOR ENGINE
# ==============================================================================
class SlideGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Theme.WIDTH
        self.prs.slide_height = Theme.HEIGHT
        self.layout_blank = self.prs.slide_layouts[6]

    def _fix_font(self, run, font_name):
        """修复字体显示（特别是中文）"""
        try:
            run.font.name = font_name
            run.font._element.rPr.set(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}ea', 
                font_name
            )
        except:
            pass

    def add_slide(self):
        slide = self.prs.slides.add_slide(self.layout_blank)
        # 设置深色背景
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = Theme.COLOR_BG_DEEP
        return slide

    def draw_text(self, slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        self._fix_font(run, Theme.FONT_CN)
        return txBox

    def draw_card(self, slide, x, y, w, h, bg_color=Theme.COLOR_BG_CARD):
        """绘制圆角卡片背景"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background() # 无边框
        shape.adjustments[0] = 0.03 # 圆角弧度
        return shape

    def draw_tag(self, slide, text, x, y, bg_color):
        """绘制胶囊标签"""
        w, h = Inches(1.2), Inches(0.4)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        shape.adjustments[0] = 1.0 # 纯圆角
        
        # 标签文字
        self.draw_text(slide, text, x, y + Inches(0.05), w, h, 12, Theme.COLOR_TEXT_WHITE, True, PP_ALIGN.CENTER)

    def draw_placeholder(self, slide, x, y, w, h, label):
        """绘制截图占位符"""
        shape = self.draw_card(slide, x, y, w, h, RGBColor(45, 50, 60))
        # 虚线边框效果 (模拟)
        self.draw_text(slide, f"[{label}]\nPLACE UI SCREENSHOT HERE", x, y + h/2 - Inches(0.3), w, Inches(1), 
                       14, Theme.COLOR_TEXT_GREY, False, PP_ALIGN.CENTER)

# ==============================================================================
# 3. CONTENT BUILDER
# ==============================================================================
def create_presentation():
    gen = SlideGenerator()

    # --------------------------------------------------------------------------
    # P1: 封面 (Cover)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    
    # 装饰性光晕 (模拟 WebGL 星球)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(8), Inches(8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(20, 40, 60)
    glow.line.fill.background()

    # 标题区
    gen.draw_text(slide, "AgriVerse", Inches(1), Inches(2.5), Inches(10), Inches(2), 
                  96, Theme.COLOR_TEXT_WHITE, True)
    
    gen.draw_text(slide, "农业产品融销平台 · 种植智慧 收获未来", Inches(1.2), Inches(4), Inches(10), Inches(1), 
                  28, Theme.COLOR_FARMER, False)

    gen.draw_text(slide, "第十小组中期答辩 | 2025年1月 | v1.2", Inches(1.2), Inches(6.5), Inches(8), Inches(0.5), 
                  16, Theme.COLOR_TEXT_GREY)

    # --------------------------------------------------------------------------
    # P2: 项目概览与技术栈 (Overview & Tech)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "项目技术全景 (Tech Stack)", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_TEXT_WHITE, True)

    # 左侧：前端架构卡片
    gen.draw_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    gen.draw_text(slide, "Frontend Core", Inches(0.8), Inches(1.8), Inches(4), Inches(0.5), 20, Theme.COLOR_TECH, True)
    
    techs = [
        ("React 18", "Concurrent Mode, Suspense"),
        ("TypeScript", "Strict Type Checking"),
        ("Vite 5.x", "Instant Server Start"),
        ("Zustand", "Atomic State Management"),
        ("Tailwind", "Utility-first CSS"),
        ("WebGL", "Three.js / React-Three-Fiber")
    ]
    
    y_pos = 2.4
    for title, desc in techs:
        gen.draw_text(slide, f"• {title}", Inches(0.8), Inches(y_pos), Inches(2), Inches(0.4), 16, Theme.COLOR_TEXT_WHITE, True)
        gen.draw_text(slide, desc, Inches(3), Inches(y_pos), Inches(3), Inches(0.4), 14, Theme.COLOR_TEXT_GREY)
        y_pos += 0.45

    # 右侧：截图展示
    gen.draw_placeholder(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5), "Home Page 3D Planet")

    # --------------------------------------------------------------------------
    # P3: 农户端核心 (Farmer Core)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "农户角色：融销一体化", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_FARMER, True)

    # 左图右文
    gen.draw_placeholder(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5), "Farmer Dashboard & Joint Loan")

    # 右侧功能列表
    x_text = Inches(8.3)
    gen.draw_card(slide, x_text, Inches(1.5), Inches(4.5), Inches(5))
    
    features = [
        ("田心星云", "数据可视化看板，实时监控产量与收益"),
        ("智能拼单 (Joint Loan)", "独创功能：金额不足自动匹配拼单，降低融资门槛"),
        ("电子签约", "全流程无纸化，在线签署法律效力合同"),
        ("田心市场", "一键发布农产品，支持批量打印电子面单")
    ]
    
    cur_y = 1.8
    for title, desc in features:
        gen.draw_text(slide, title, x_text + Inches(0.2), Inches(cur_y), Inches(4), Inches(0.4), 18, Theme.COLOR_FARMER, True)
        gen.draw_text(slide, desc, x_text + Inches(0.2), Inches(cur_y + 0.35), Inches(4), Inches(0.6), 14, Theme.COLOR_TEXT_GREY)
        cur_y += 1.1

    # --------------------------------------------------------------------------
    # P4: 核心亮点：智能拼单逻辑 (Joint Loan Logic)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "核心业务逻辑：智能拼单 (Joint Loan)", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_FARMER, True)

    # 流程图模拟
    card_w = Inches(3.5)
    card_h = Inches(4)
    gap = Inches(0.5)
    start_x = Inches(1)
    
    # Step 1
    gen.draw_card(slide, start_x, Inches(2), card_w, card_h)
    gen.draw_text(slide, "01. 发起申请", start_x + Inches(0.2), Inches(2.2), card_w, Inches(0.5), 20, Theme.COLOR_TEXT_WHITE, True)
    gen.draw_text(slide, "用户申请额度 < 银行最低门槛\n系统自动触发“拼单模式”", start_x + Inches(0.2), Inches(3), card_w - Inches(0.4), Inches(2), 16, Theme.COLOR_TEXT_GREY)

    # Step 2
    gen.draw_card(slide, start_x + card_w + gap, Inches(2), card_w, card_h)
    gen.draw_text(slide, "02. 智能匹配", start_x + card_w + gap + Inches(0.2), Inches(2.2), card_w, Inches(0.5), 20, Theme.COLOR_TEXT_WHITE, True)
    gen.draw_text(slide, "算法根据：\n- 信用评分\n- 资金需求量\n- 种植周期\n自动推荐最佳拼单伙伴", start_x + card_w + gap + Inches(0.2), Inches(3), card_w - Inches(0.4), Inches(2), 16, Theme.COLOR_TEXT_GREY)

    # Step 3
    gen.draw_card(slide, start_x + (card_w + gap) * 2, Inches(2), card_w, card_h)
    gen.draw_text(slide, "03. 联合授信", start_x + (card_w + gap) * 2 + Inches(0.2), Inches(2.2), card_w, Inches(0.5), 20, Theme.COLOR_TEXT_WHITE, True)
    gen.draw_text(slide, "拼单组队成功 -> 形成资产包\n银行统一授信 -> 分别放款", start_x + (card_w + gap) * 2 + Inches(0.2), Inches(3), card_w - Inches(0.4), Inches(2), 16, Theme.COLOR_TEXT_GREY)

    # --------------------------------------------------------------------------
    # P5: 买家端 (Buyer Role)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "买家角色：沉浸式采购体验", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_BUYER, True)

    # 上半部分：Dashboard
    gen.draw_placeholder(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(5), "Buyer Dashboard")
    
    # 中间：Shopping Cart
    gen.draw_placeholder(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(5), "Cart & Product Detail")
    
    # 右侧：功能说明
    gen.draw_card(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(5))
    gen.draw_text(slide, "主要功能", Inches(9.1), Inches(1.8), Inches(3), Inches(0.5), 20, Theme.COLOR_BUYER, True)
    buyer_feats = ["购市星云 (采购驾驶舱)", "多维度商品对比", "购物车与分期支付", "全流程退款追踪"]
    by = 2.5
    for f in buyer_feats:
        gen.draw_text(slide, f"• {f}", Inches(9.1), Inches(by), Inches(3.5), Inches(0.5), 16, Theme.COLOR_TEXT_WHITE)
        by += 0.6

    # --------------------------------------------------------------------------
    # P6: 银行端 (Bank Role)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "银行角色：智能风控体系", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_BANK, True)

    # 大屏风控看板占位
    gen.draw_placeholder(slide, Inches(0.5), Inches(1.5), Inches(8.5), Inches(5), "Risk Control Cockpit (Charts)")

    # 右侧面板
    gen.draw_card(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(5))
    gen.draw_text(slide, "信贷工厂", Inches(9.5), Inches(1.8), Inches(3), Inches(0.5), 20, Theme.COLOR_BANK, True)
    
    bank_feats = [("评分卡模型", "A/B/C卡自动评分"), ("贷后预警", "资金流向异常监控"), ("电子合同", "一键生成/签署"), ("放款管理", "自动化放款流程")]
    by = 2.5
    for t, d in bank_feats:
        gen.draw_text(slide, t, Inches(9.5), Inches(by), Inches(3), Inches(0.3), 16, Theme.COLOR_TEXT_WHITE, True)
        gen.draw_text(slide, d, Inches(9.5), Inches(by + 0.25), Inches(3), Inches(0.3), 12, Theme.COLOR_TEXT_GREY)
        by += 0.9

    # --------------------------------------------------------------------------
    # P7: 专家与管理员 (Expert & Admin)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "生态支撑：专家服务与运营管理", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_TEXT_WHITE, True)

    # 左半边：专家
    gen.draw_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    gen.draw_text(slide, "🎓 专家 (Expert)", Inches(0.8), Inches(1.8), Inches(4), Inches(0.5), 24, Theme.COLOR_EXPERT, True)
    gen.draw_placeholder(slide, Inches(0.8), Inches(2.5), Inches(5.4), Inches(3), "Expert Q&A / Knowledge")
    gen.draw_text(slide, "• 付费问答与知识变现\n• 预约咨询日历管理", Inches(0.8), Inches(5.8), Inches(5), Inches(1), 16, Theme.COLOR_TEXT_GREY)

    # 右半边：管理员
    gen.draw_card(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5))
    gen.draw_text(slide, "🛡️ 管理员 (Admin)", Inches(7.1), Inches(1.8), Inches(4), Inches(0.5), 24, Theme.COLOR_ADMIN, True)
    gen.draw_placeholder(slide, Inches(7.1), Inches(2.5), Inches(5.4), Inches(3), "Admin Operation Center")
    gen.draw_text(slide, "• 内容/商品/专家三审机制\n• 全平台权限与日志监控", Inches(7.1), Inches(5.8), Inches(5), Inches(1), 16, Theme.COLOR_TEXT_GREY)

    # --------------------------------------------------------------------------
    # P8: 导航与架构 (Navigation Architecture)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "技术亮点：三层导航架构", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_TECH, True)

    # 架构图示
    base_y = 2.0
    
    # Layer 1
    gen.draw_card(slide, Inches(2), Inches(base_y), Inches(9), Inches(1), Theme.COLOR_TECH)
    gen.draw_text(slide, "Level 1: 顶部主导航 (Tab Switching)", Inches(2.5), Inches(base_y + 0.2), Inches(8), Inches(0.5), 20, RGBColor(0,0,0), True, PP_ALIGN.CENTER)
    
    # Arrow
    gen.draw_text(slide, "⬇ navigateToTab()", Inches(0), Inches(base_y + 1), Theme.WIDTH, Inches(0.5), 14, Theme.COLOR_TEXT_GREY, False, PP_ALIGN.CENTER)

    # Layer 2
    gen.draw_card(slide, Inches(3), Inches(base_y + 1.5), Inches(7), Inches(1), RGBColor(0, 150, 136))
    gen.draw_text(slide, "Level 2: 页面子路由 (Sub-Route System)", Inches(3.5), Inches(base_y + 1.7), Inches(6), Inches(0.5), 20, Theme.COLOR_TEXT_WHITE, True, PP_ALIGN.CENTER)
    
    # Arrow
    gen.draw_text(slide, "⬇ navigateToSubRoute()", Inches(0), Inches(base_y + 2.5), Theme.WIDTH, Inches(0.5), 14, Theme.COLOR_TEXT_GREY, False, PP_ALIGN.CENTER)

    # Layer 3
    gen.draw_card(slide, Inches(4), Inches(base_y + 3.0), Inches(5), Inches(1), RGBColor(0, 100, 100))
    gen.draw_text(slide, "Level 3: 移动端底部导航 (Mobile Only)", Inches(4.5), Inches(base_y + 3.2), Inches(4), Inches(0.5), 20, Theme.COLOR_TEXT_WHITE, True, PP_ALIGN.CENTER)

    # --------------------------------------------------------------------------
    # P9: 状态管理 (Zustand)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "技术亮点：Zustand 状态管理", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_TECH, True)
    
    gen.draw_placeholder(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(5), "Store Code Snippet")
    
    gen.draw_card(slide, Inches(6), Inches(1.5), Inches(6.8), Inches(5))
    gen.draw_text(slide, "模块化 Store 设计", Inches(6.2), Inches(1.8), Inches(6), Inches(0.5), 24, Theme.COLOR_TEXT_WHITE, True)
    
    stores = [
        "financingStore: 封装利息计算、还款计划生成逻辑",
        "cartStore: 本地持久化、金额自动汇总",
        "userStore: 角色权限控制 (RBAC)、Token管理",
        "msgStore: 全局消息通知、WebSocket 状态同步"
    ]
    
    sy = 2.5
    for s in stores:
        gen.draw_text(slide, f"• {s}", Inches(6.2), Inches(sy), Inches(6.4), Inches(0.5), 18, Theme.COLOR_TEXT_GREY)
        sy += 0.8

    # --------------------------------------------------------------------------
    # P10: 项目完成度 (Project Status)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "项目完成度 (Status Report)", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_TEXT_WHITE, True)

    # 3个大数据卡片
    metrics = [
        ("100%", "API 接口联调", "115/115 接口全通", Theme.COLOR_FARMER),
        ("60+", "功能页面", "覆盖5大角色全流程", Theme.COLOR_TECH),
        ("V1.2", "当前版本", "架构稳定 交互流畅", Theme.COLOR_ADMIN)
    ]
    
    mx = Inches(0.5)
    mw = Inches(3.8)
    for num, title, sub, color in metrics:
        gen.draw_card(slide, mx, Inches(2), mw, Inches(4))
        gen.draw_text(slide, num, mx, Inches(2.5), mw, Inches(1.5), 80, color, True, PP_ALIGN.CENTER)
        gen.draw_text(slide, title, mx, Inches(4), mw, Inches(0.5), 20, Theme.COLOR_TEXT_WHITE, True, PP_ALIGN.CENTER)
        gen.draw_text(slide, sub, mx, Inches(4.6), mw, Inches(0.5), 16, Theme.COLOR_TEXT_GREY, False, PP_ALIGN.CENTER)
        mx += mw + Inches(0.4)

    # --------------------------------------------------------------------------
    # P11: 未来展望 (Future)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "未来计划 (Future Roadmap)", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), 32, Theme.COLOR_TEXT_WHITE, True)
    
    gen.draw_card(slide, Inches(1), Inches(2), Inches(11.3), Inches(4))
    
    timeline = [
        ("Phase 1: 测试", "E2E 自动化测试覆盖核心流程 (Cypress)"),
        ("Phase 2: 性能", "微前端物理拆分，路由懒加载优化"),
        ("Phase 3: 生态", "PWA 离线支持，i18n 多语言适配"),
        ("Phase 4: 实时", "WebSocket 消息全链路打通")
    ]
    
    ty = 2.4
    for title, desc in timeline:
        gen.draw_text(slide, title, Inches(1.5), Inches(ty), Inches(4), Inches(0.5), 20, Theme.COLOR_TECH, True)
        gen.draw_text(slide, desc, Inches(5), Inches(ty), Inches(7), Inches(0.5), 20, Theme.COLOR_TEXT_GREY)
        ty += 0.8

    # --------------------------------------------------------------------------
    # P12: 结束页 (End)
    # --------------------------------------------------------------------------
    slide = gen.add_slide()
    gen.draw_text(slide, "AgriVerse", Inches(0), Inches(3), Theme.WIDTH, Inches(1.5), 60, Theme.COLOR_TEXT_WHITE, True, PP_ALIGN.CENTER)
    gen.draw_text(slide, "感谢观看 · Q&A", Inches(0), Inches(4.5), Theme.WIDTH, Inches(1), 24, Theme.COLOR_FARMER, False, PP_ALIGN.CENTER)

    # Save
    gen.prs.save("AgriVerse_Ultimate.pptx")
    print("✅ AgriVerse_Ultimate.pptx 生成成功！(共12页)")

if __name__ == "__main__":
    create_presentation()